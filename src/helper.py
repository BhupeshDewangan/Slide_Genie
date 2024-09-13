import base64
from pptx import Presentation
from pptx.util import Inches, Pt
import json
import google.generativeai as genai
import os
import streamlit as st
from io import BytesIO
from urllib.request import urlopen


# API KEY
genai.configure(api_key=st.secrets["api_key"])
model = genai.GenerativeModel(model_name="gemini-1.5-flash")

# Sample Type
f = open('JSON/Sample.json')
data = json.load(f)


def get_synthetic_data(number, subject, language, tone):
    """
    Create synthetic data for creating the PowerPoints in JSON format.

  Args:
    The number of PowerPoint files
    The subject for these files
    """

    prompt = f"""I am automating the creation of PowerPoint files using JSON data.
    Here is a sample JSON file that contains information for two PowerPoint files.
    "AI in action" and "AI in action": {data}
    In the same manner, can you create information for {number} PowerPoint files
    into a single JSON file format on the subject {subject}?
    Please replace the 'lorem ipsum' content with actual generated content and 
    the values of JSON will be appropriate names in {language}.
    The JSON File should not have the {subject} name anywhere.
    Please use the same image links in the source JSON file.
    Please provide only the code, don't include any explanations in your responses.
    'first_column' and 'second_column' should have at least 200 words of text.
    Also, don't use any escape characters in the JSON.
    The tone of the content should be {tone}."""


    response = model.generate_content(prompt)
    spe_char_to_remove = ['`', '`', 'json']
    output =  response.text
    for character in spe_char_to_remove:
      output = output.replace(character, '')

    return output



def change_to_list(item, str): 
  """
  Identify if the string has list content and creates text content based on that

  Args:
    The content placeholder
    The content
  """
  tf = item.text_frame
  if "\n" in str:
      list_array = str.split("\n")    
      for i in list_array:
          print(i)
          p = tf.add_paragraph()
          p.text = i        
          
  else:
      tf.text = str



def create_presentation(prs, presentation_data):
  """
  Creates a slide in the presentation based on the layout and content provided.

  Args:
    presentation_data: A dictionary containing the slide layout and content.
    prs: The presentation object.
  """
  layout_mapping = {
      "single_title": 0,
      "double-column_text": 1,
      "text_and_image": 2,
      "image_and_text": 3,
      "image_quote": 4,
      "qa": 5,
      "thankyou": 6
  }

  for layout_name, content in presentation_data.items():
    layout_index = layout_mapping[layout_name]
    slide_layout = prs.slide_layouts[layout_index]    
    slide = prs.slides.add_slide(slide_layout)  

    if layout_name == "single_title":
      slide.placeholders[0].text = content["title"]  
      slide.placeholders[1].text = content["subtitle"] 
      url = (content["image"])
      image_data = BytesIO(urlopen(url).read())      
      picture = slide.placeholders[13].insert_picture(image_data) 

    elif layout_name == "double-column_text":            
      slide.placeholders[0].text = content["title"]                
      change_to_list(slide.placeholders[1], content["first_column"])
      change_to_list(slide.placeholders[13], content["second_column"])   

    elif layout_name == "text_and_image":
      slide.placeholders[0].text = content["title"]      
      change_to_list(slide.placeholders[1], content["first_column"])       
      url = (content["image"])
      image_data = BytesIO(urlopen(url).read())      
      picture = slide.placeholders[13].insert_picture(image_data)

    elif layout_name == "image_and_text":
      slide.placeholders[0].text = content["title"]      
      change_to_list(slide.placeholders[1], content["first_column"])
      url = (content["image"])
      image_data = BytesIO(urlopen(url).read())      
      picture = slide.placeholders[13].insert_picture(image_data)

    elif layout_name == "image_quote":
      slide.placeholders[0].text = content["title"]      
      change_to_list(slide.placeholders[14], content["quote"])
      url = (content["image"])
      image_data = BytesIO(urlopen(url).read())      
      picture = slide.placeholders[13].insert_picture(image_data)

    elif layout_name == "qa":
      slide.placeholders[0].text = content["title"] 

    elif layout_name == "thankyou":
      # slide.placeholders[0].text = content["title"]    
      change_to_list(slide.placeholders[14], content["wide_text"]) 



def automate_ppts(text, template):
  """
  Creates a number of PPTs based on the data json.
  """

  # print(template)
  output_dir = f"PPTs/{text}"
  os.makedirs(output_dir, exist_ok=True)

  with open("JSON/presentation_data.json", "r", encoding='utf-8') as f:
    presentation_data = json.load(f)
 
  for name, content in presentation_data.items():       
    prs = Presentation(template)
    create_presentation(prs, content)
    # prs.save(f"PPTs/{text}/{name}.pptx")
    prs.save(os.path.join(output_dir, f"{name}.pptx"))



def show_pdf(file_path):
    with open(file_path, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')
    # pdf_display = f'<embed src="data:application/pdf;base64,{base64_pdf}" width="700" height="1000" type="application/pdf">'
    pdf_display = F'<iframe src="data:application/pdf;base64,{base64_pdf}" width="700" height="1000" type="application/pdf"></iframe>'
    st.markdown(pdf_display, unsafe_allow_html=True)